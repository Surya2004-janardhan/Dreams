from flask import Flask, request, send_file, after_this_request
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
import io
import os

app = Flask(__name__)

@app.route('/generate', methods=['POST'])
def generate():
    data = request.get_json()
    title = data['title']
    content = data['content']
    
    # Create image
    img = Image.open('assets/Post-Base-Image.png')
    draw = ImageDraw.Draw(img)
    try:
        font_title = ImageFont.truetype("arialbd.ttf", 48)  # Bold Arial for title
    except:
        font_title = ImageFont.truetype("arial.ttf", 48)
    try:
        font_content = ImageFont.truetype("arial.ttf", 24)  # Regular Arial for content
    except:
        font_content = ImageFont.load_default()
    
    # Margins
    margin = 50
    
    # Draw title
    draw.text((margin, margin), title, fill="black", font=font_title)
    
    # Draw content
    draw.text((margin, margin + 120), content, fill=(100, 100, 100), font=font_content)
    
    # Save to assets
    img_path = 'assets/Post-Image.png'
    img.save(img_path)
    
    # Create slide
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content
    prs.save('slide/generated_slide.pptx')
    
    # Return image
    @after_this_request
    def remove_file(response):
        try:
            os.remove(img_path)
        except Exception as e:
            app.logger.error(f"Error deleting file: {e}")
        return response
    
    return send_file(img_path, mimetype='image/png', as_attachment=True, download_name='Post-Image.png')

if __name__ == '__main__':
    app.run(debug=True)